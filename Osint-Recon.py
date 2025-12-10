# Basic imports and configuration
import streamlit as st
import streamlit.components.v1 as components
from pyvis.network import Network
import requests
import whois
import socket
import pandas as pd
import concurrent.futures
import tempfile
import uuid
import json
import os
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches, Pt
from datetime import datetime

# Page setup and styling
st.set_page_config(page_title="OSINT-Recon", layout="wide", page_icon="🌐")
st.markdown(
    """
    <style>
    .stApp { background-color: #0E1117; color: #E0E0E0; }
    .metric-box { background-color: #262730; padding: 15px; border-radius: 10px; border-left: 5px solid #00ADB5; }
    .small-muted { color: #9aa0a6; font-size:12px; }
    </style>
    """,
    unsafe_allow_html=True,
)

# HTTP headers and file paths
HEADERS = {
    "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
                  "AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36"
}
REPUTATION_FILE = "demo_reputation.json"
DEMO_DIR = "demo"

# Create demo reputation files if missing
if not os.path.exists(DEMO_DIR):
    os.makedirs(DEMO_DIR, exist_ok=True)
if not os.path.exists(REPUTATION_FILE):
    demo_reputation = {
        "bad_ips": ["203.0.113.10", "198.51.100.5"],
        "bad_domains": ["malicious-example.com", "badactor.test"]
    }
    with open(REPUTATION_FILE, "w") as f:
        json.dump(demo_reputation, f, indent=2)

# Helper functions for safe extraction and loading reputation data
def safe_first(val):
    if isinstance(val, (list, tuple, set)):
        return next(iter(val), None)
    return val

def load_reputation():
    try:
        with open(REPUTATION_FILE, "r") as f:
            return json.load(f)
    except Exception:
        return {"bad_ips": [], "bad_domains": []}

REPUTATION = load_reputation()

# Subdomain enumeration - crt.sh lookup
@st.cache_data(show_spinner=False)
def get_subdomains(domain):
    url = f"https://crt.sh/?q=%.{domain}&output=json"
    try:
        resp = requests.get(url, headers=HEADERS, timeout=12)
        if resp.status_code != 200:
            return []
        data = resp.json()
        subdomains = set()
        for entry in data:
            name = entry.get("name_value", "")
            for n in str(name).splitlines():
                n = n.strip().lower()
                if n and "*" not in n:
                    subdomains.add(n)
        return sorted(subdomains)
    except Exception:
        return []

# IP geolocation lookup using ip-api
@st.cache_data(show_spinner=False)
def get_ip_location(ip_address):
    url = f"http://ip-api.com/json/{ip_address}?fields=status,message,country,countryCode,regionName,city,zip,lat,lon,timezone,isp,org,as,query"
    try:
        resp = requests.get(url, timeout=6)
        return resp.json()
    except Exception:
        return {"status": "fail", "message": "API error"}

# WHOIS lookup for domain metadata
@st.cache_data(show_spinner=False)
def get_whois_data(domain):
    try:
        w = whois.whois(domain)
        registrar = safe_first(getattr(w, "registrar", None)) or "Hidden"
        creation = safe_first(getattr(w, "creation_date", None))
        creation_str = str(creation) if creation else "Hidden"
        emails = safe_first(getattr(w, "emails", None)) or "Hidden"
        return {"Registrar": registrar, "Creation Date": creation_str, "Emails": emails}
    except Exception:
        return {"Error": "WHOIS Hidden/Failed"}

# Build interactive graph (PyVis)
def create_graph_html(target, nodes):
    net = Network(height="450px", width="100%", bgcolor="#111111", font_color="white")
    net.toggle_physics(False)
    net.add_node(target, label=target, color="#FF4B4B", size=30, title="Target")
    for n in nodes[:40]:
        net.add_node(n, label=n, color="#00ADB5", size=16, title="Subdomain")
        net.add_edge(target, n)
    tmpname = f"network_{uuid.uuid4().hex}.html"
    tmp_path = os.path.join(tempfile.gettempdir(), tmpname)
    net.save_graph(tmp_path)
    return tmp_path

# Common ports list for scanning
COMMON_PORTS = {
    21: "FTP", 22: "SSH", 23: "Telnet", 25: "SMTP",
    53: "DNS", 80: "HTTP", 443: "HTTPS", 445: "SMB",
    3306: "MySQL", 3389: "RDP", 5432: "PostgreSQL", 8080: "HTTP-Proxy"
}

# Single-port connect scan
def scan_port_connect(ip, port, timeout=0.6):
    try:
        with socket.create_connection((ip, port), timeout=timeout):
            return port
    except Exception:
        return None

# Threaded port scanning
def scan_ports_threaded(ip, ports=COMMON_PORTS, max_workers=40):
    open_ports = []
    results = {}
    with concurrent.futures.ThreadPoolExecutor(max_workers=max_workers) as ex:
        futures = {ex.submit(scan_port_connect, ip, p): p for p in ports}
        for fut in concurrent.futures.as_completed(futures):
            p = futures[fut]
            res = fut.result()
            if res:
                open_ports.append(p)
    for p, name in ports.items():
        results[p] = {"Service": name, "Status": ("OPEN 🟢" if p in open_ports else "Closed 🔴")}
    return results

# PTR (reverse DNS) lookup
def get_reverse_dns(ip):
    try:
        return socket.gethostbyaddr(ip)[0]
    except Exception:
        return "No PTR Record Found"

# PDF report generator
def generate_pdf(report_type, target, data_dict, subdomains):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", 'B', 16)
    pdf.cell(0, 10, txt=f"{report_type} Report: {target}", ln=True, align='C')
    pdf.ln(6)
    pdf.set_font("Arial", size=11)
    pdf.cell(0, 8, txt=f"Generated: {datetime.utcnow().isoformat()} UTC", ln=True)
    pdf.ln(6)
    pdf.set_font("Arial", 'B', 12)
    pdf.cell(0, 8, txt="Extracted Intelligence:", ln=True)
    pdf.set_font("Arial", size=10)
    pdf.ln(4)
    if isinstance(data_dict, dict):
        for k, v in data_dict.items():
            pdf.multi_cell(0, 6, txt=f"{k}: {v}")
    if isinstance(subdomains, (list, tuple)):
        pdf.ln(4)
        pdf.set_font("Arial", 'B', 11)
        pdf.cell(0, 8, txt=f"Subdomains Found ({len(subdomains)}):", ln=True)
        pdf.set_font("Arial", size=9)
        for item in subdomains[:200]:
            pdf.multi_cell(0, 5, txt=f"- {item}")
    out_name = f"Recon_Report_{report_type}_{uuid.uuid4().hex[:6]}.pdf"
    pdf.output(out_name)
    return out_name

# PPTX slide generator
def generate_pptx(report_type, target, metrics: dict, notes: str):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(1))
    tf = title.text_frame
    tf.text = f"{report_type} Summary: {target}"
    tf.paragraphs[0].font.size = Pt(24)
    box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(4.5), Inches(3.5))
    bf = box.text_frame
    bf.word_wrap = True
    bf.text = "Key Metrics\n"
    for k, v in metrics.items():
        p = bf.add_paragraph()
        p.text = f"- {k}: {v}"
        p.level = 1
    note_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.8), Inches(4), Inches(3.5))
    nf = note_box.text_frame
    nf.text = "Notes\n"
    p = nf.add_paragraph()
    p.text = notes
    p.level = 1
    out_name = f"Recon_Summary_{report_type}_{uuid.uuid4().hex[:6]}.pptx"
    prs.save(out_name)
    return out_name

# Session state helpers
if "running" not in st.session_state:
    st.session_state.running = False

def set_running(val: bool):
    st.session_state.running = val

# Main UI rendering
def main():
    st.title("🌐 OSINT ReconSuite V3 — Improved")
    st.markdown(
        "Domain recon | IP intelligence | Safe, demo-friendly outputs.",
        unsafe_allow_html=True,
    )
    st.info("Use only on authorized targets.", icon="⚖️")

    tab_domain, tab_ip, tab_demo = st.tabs(["🔎 Domain Recon", "📍 IP Intel", "📁 Demo & Files"])

    # Domain recon UI
    with tab_domain:
        st.subheader("Domain Surface Attack Analysis")
        domain_input = st.text_input("Target Domain", "tesla.com")
        col1, col2 = st.columns([1, 2])
        with col1:
            start_btn = st.button("Start Domain Scan", disabled=st.session_state.running)

        if start_btn:
            try:
                set_running(True)
                with st.spinner("Mapping Attack Surface..."):
                    subs = get_subdomains(domain_input)
                    whois_info = get_whois_data(domain_input)
                    c1, c2, c3 = st.columns(3)
                    c1.metric("Subdomains Found", len(subs))
                    c2.metric("Registrar", whois_info.get("Registrar", "Unknown"))
                    c3.metric("Creation Date", whois_info.get("Creation Date", "Unknown"))

                    domain_flag = "Suspicious" if domain_input.lower() in REPUTATION.get("bad_domains", []) else "Clean"
                    st.markdown(f"**Local Reputation:** `{domain_flag}`")

                    st.subheader("🔗 Attack Surface Graph")
                    if subs:
                        graph_path = create_graph_html(domain_input, subs)
                        with open(graph_path, 'r', encoding='utf-8') as f:
                            html = f.read()
                        components.html(html, height=480)
                        st.write(f"Graph saved: `{graph_path}`")
                    else:
                        st.warning("No subdomains found.")

                    st.write("### WHOIS Data")
                    st.json(whois_info)

                    with st.expander(f"View All Subdomains ({len(subs)})"):
                        st.write(subs)

                    pdf_name = generate_pdf("Domain", domain_input, whois_info, subs)
                    ppt_name = generate_pptx("Domain", domain_input,
                                             {"Subdomains": len(subs), "Registrar": whois_info.get("Registrar", "Unknown")},
                                             "Auto-generated summary.")
                    with open(pdf_name, "rb") as f:
                        st.download_button("Download Domain Report (PDF)", f, file_name=pdf_name)
                    with open(ppt_name, "rb") as f:
                        st.download_button("Download Domain Summary (PPTX)", f, file_name=ppt_name)
            finally:
                set_running(False)

    # IP intelligence UI
    with tab_ip:
        st.subheader("IP Infrastructure & Geo-Location")
        ip_input = st.text_input("Enter IP Address", "8.8.8.8")
        col1, col2 = st.columns([1, 2])
        with col1:
            start_ip = st.button("Scan IP Address", disabled=st.session_state.running)

        if start_ip:
            try:
                set_running(True)
                with st.spinner("Processing..."):
                    loc_data = get_ip_location(ip_input)
                    if loc_data.get("status") == "success":
                        ptr = get_reverse_dns(ip_input)
                        port_results = scan_ports_threaded(ip_input)
                        m1, m2, m3 = st.columns(3)
                        m1.metric("ISP / Org", loc_data.get("isp", "Unknown"))
                        m2.metric("Location", f"{loc_data.get('city','')}, {loc_data.get('country','')}")
                        m3.metric("Reverse DNS", ptr)

                        st.write("### Physical Location")
                        try:
                            df_map = pd.DataFrame({'lat': [loc_data.get('lat')], 'lon': [loc_data.get('lon')]})
                            st.map(df_map, zoom=4)
                        except:
                            st.write("Map unavailable.")

                        colA, colB = st.columns(2)
                        with colA:
                            st.write("### IP Details")
                            details = {
                                "IP": loc_data.get("query"),
                                "Region": loc_data.get("regionName"),
                                "Zip Code": loc_data.get("zip"),
                                "Timezone": loc_data.get("timezone"),
                                "AS Number": loc_data.get("as"),
                                "Org": loc_data.get("org")
                            }
                            st.table(pd.DataFrame(list(details.items()), columns=["Field", "Value"]))

                        with colB:
                            st.write("### Port Scan Results")
                            df_ports = pd.DataFrame.from_dict(port_results, orient='index')
                            st.dataframe(df_ports, use_container_width=True)
                            open_count = sum(1 for v in port_results.values() if "OPEN" in v["Status"])
                            if open_count > 0:
                                st.warning(f"Found {open_count} open ports.")
                            else:
                                st.success("No common ports exposed.")

                        ip_flag = "Suspicious" if ip_input in REPUTATION.get("bad_ips", []) else "Clean"
                        st.markdown(f"**Local Reputation:** `{ip_flag}`")

                        pdf_name = generate_pdf("IP_Intel", ip_input, details, [])
                        ppt_name = generate_pptx("IP", ip_input,
                                                 {"Open Ports": open_count, "ISP": loc_data.get("isp", "")},
                                                 "Auto-generated summary.")
                        with open(pdf_name, "rb") as f:
                            st.download_button("Download IP Report (PDF)", f, file_name=pdf_name)
                        with open(ppt_name, "rb") as f:
                            st.download_button("Download IP Summary (PPTX)", f, file_name=ppt_name)
                    else:
                        st.error("Failed to fetch IP details.")
            finally:
                set_running(False)

    # Demo & sample file utilities
    with tab_demo:
        st.subheader("Demo & Supporting Files")
        st.markdown("- Demo reputation file.\n- Works offline for testing.")

        if st.button("Show demo reputation file"):
            st.code(json.dumps(REPUTATION, indent=2), language="json")

        st.markdown("### Create demo sample files")
        if st.button("Create demo domain + IP samples"):
            demo_subs = ["test.tesla.com", "dev.tesla.com", "assets.tesla.com"]
            with open(os.path.join(DEMO_DIR, "sample_subdomains.txt"), "w") as f:
                f.write("\n".join(demo_subs))
            with open(os.path.join(DEMO_DIR, "sample_ip.json"), "w") as f:
                json.dump({"ip": "8.8.8.8", "sample": True}, f, indent=2)
            st.success(f"Demo files created in `{DEMO_DIR}/`")

        st.markdown("### Notes")
        st.markdown("- Uses passive OSINT.\n- Respect API limits.\n- Only scan authorized systems.")

if __name__ == "__main__":
    main()
